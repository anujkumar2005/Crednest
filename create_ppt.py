"""
CredNest AI - PowerPoint Presentation Generator
Converts the Markdown presentation to PowerPoint format
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def create_presentation():
    """Create PowerPoint presentation from markdown content"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (Dark theme)
    DARK_BG = RGBColor(10, 10, 10)  # #0a0a0a
    LIGHT_BG = RGBColor(26, 26, 26)  # #1a1a1a
    ACCENT_GOLD = RGBColor(255, 215, 0)  # #FFD700
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_GRAY = RGBColor(200, 200, 200)
    
    # Read markdown content
    with open('CredNest_AI_Presentation.md', 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split into slides
    slides_content = content.split('---\n\n## Slide ')
    
    for slide_num, slide_text in enumerate(slides_content[1:], 1):  # Skip first split
        # Parse slide content
        lines = slide_text.strip().split('\n')
        title_line = lines[0].split(': ', 1)
        
        if len(title_line) > 1:
            slide_number = title_line[0]
            slide_title = title_line[1]
        else:
            slide_title = lines[0]
        
        # Create slide
        if slide_num == 1:
            # Title slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            add_title_slide(slide, DARK_BG, ACCENT_GOLD, TEXT_WHITE)
        elif 'Table of Contents' in slide_title:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_toc_slide(slide, lines[2:], LIGHT_BG, ACCENT_GOLD, TEXT_WHITE)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_content_slide(slide, slide_title, lines[2:], LIGHT_BG, ACCENT_GOLD, TEXT_WHITE, TEXT_GRAY)
    
    # Save presentation
    prs.save('CredNest_AI_Presentation.pptx')
    print("✅ PowerPoint presentation created successfully!")
    print("📄 File: CredNest_AI_Presentation.pptx")

def add_title_slide(slide, bg_color, accent_color, text_color):
    """Add title slide with student details"""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "CredNest AI"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = accent_color
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Financial Management Platform"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.italic = True
    subtitle_para.font.color.rgb = text_color
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Student details
    details_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(2))
    details_frame = details_box.text_frame
    
    details = [
        "Presented by:",
        "Anujkumar",
        "",
        "Register Number: 22ETAI410010",
        "",
        "M S Ramaiah University of Applied Science",
        "Department of Computer Science & Engineering",
        "",
        "November 2025"
    ]
    
    for i, detail in enumerate(details):
        if i > 0:
            details_frame.add_paragraph()
        para = details_frame.paragraphs[i]
        para.text = detail
        para.font.size = Pt(18) if i in [0, 1] else Pt(14)
        para.font.bold = i in [1]
        para.font.color.rgb = accent_color if i in [1] else text_color
        para.alignment = PP_ALIGN.CENTER

def add_toc_slide(slide, content_lines, bg_color, accent_color, text_color):
    """Add table of contents slide"""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Table of Contents"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = accent_color
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
    content_frame = content_box.text_frame
    
    for i, line in enumerate(content_lines):
        if line.strip():
            if i > 0:
                content_frame.add_paragraph()
            para = content_frame.paragraphs[-1]
            para.text = line.strip()
            para.font.size = Pt(16)
            para.font.color.rgb = text_color
            para.space_before = Pt(6)

def add_content_slide(slide, title, content_lines, bg_color, accent_color, text_color, gray_color):
    """Add content slide"""
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = accent_color
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if line.strip():
            if i > 0:
                content_frame.add_paragraph()
            para = content_frame.paragraphs[-1]
            
            # Remove markdown formatting
            clean_line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
            clean_line = re.sub(r'\*(.*?)\*', r'\1', clean_line)
            clean_line = re.sub(r'`(.*?)`', r'\1', clean_line)
            
            para.text = clean_line.strip()
            
            # Styling based on content
            if line.startswith('###'):
                para.font.size = Pt(24)
                para.font.bold = True
                para.font.color.rgb = accent_color
            elif line.startswith('**'):
                para.font.size = Pt(18)
                para.font.bold = True
                para.font.color.rgb = text_color
            elif line.startswith('- ') or line.startswith('* '):
                para.font.size = Pt(14)
                para.font.color.rgb = gray_color
                para.level = 1
            else:
                para.font.size = Pt(16)
                para.font.color.rgb = text_color

if __name__ == '__main__':
    try:
        create_presentation()
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        print("\n💡 Make sure python-pptx is installed:")
        print("   pip install python-pptx")
